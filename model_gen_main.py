import os
import re
import pdfkit
import requests
import subprocess
from pathlib import Path
from dotenv import load_dotenv
from googlesearch import search
from bs4 import BeautifulSoup as bs
import arxiv
import time

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Cm, Pt
from datetime import datetime


import tiktoken
from functools import partial

# LlamaIndex
from llama_index.readers.file import PDFReader
from llama_index.core.retrievers import VectorIndexRetriever
from llama_index.core.postprocessor import SimilarityPostprocessor
from llama_index.core import (
    Settings,
    get_response_synthesizer,
    StorageContext,
    VectorStoreIndex,
    load_index_from_storage,
)
from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
from openai import AzureOpenAI
from llama_index.core.query_engine import RetrieverQueryEngine
from llama_index.llms import azure_openai
import openai
from pinecone import Pinecone
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.azure_openai import AzureOpenAIEmbedding
from llama_index.vector_stores.pinecone import PineconeVectorStore
from llama_index.core.postprocessor import SimilarityPostprocessor, LLMRerank

load_dotenv()


class ContentGenerator:
    def __init__(self) -> None:
        self.dirkom = dirkom
        self.deskripsi_dirkom = deskripsi_dirkom
        self.pc = Pinecone(api_key=os.environ.get("PINECONE_API_KEY"))
        self.pc_index = self.pc.Index(host=os.environ.get("PINECONE_HOST"))

    def request_generate(self, prompt):
        Settings.embed_model = AzureOpenAIEmbedding(
            model="text-embedding-3-large",
            deployment_name="corpu-text-embedding-3-large",
            api_key=os.getenv("AZURE_API_KEY"),
            azure_endpoint=os.getenv("AZURE_API_BASE"),
            api_version="2023-05-15",
        )

        Settings.llm = azure_openai.AzureOpenAI(
            model="text-davinci-003",
            deployment_name="corpu-text-davinci-003",
            temperature=0.3,
            api_key=os.getenv("AZURE_API_KEY"),
            azure_endpoint=os.getenv("AZURE_API_BASE"),
            api_version="2023-05-15",
        )
        vector_store = PineconeVectorStore(self.pc_index)

        index = VectorStoreIndex.from_vector_store(vector_store)

        retriever = VectorIndexRetriever(index=index, similarity_top_k=5)

        query_engine = RetrieverQueryEngine.from_args(
            retriever=retriever,
        )

        result = query_engine.query(prompt)

        return result

    def prompt_generate_title(self, dirkom, kompetensi):
        # prompt = f"""Anda adalah seorang ahli {dirkom} yang sedang memberikan pelatihan kepada sekelompok orang. kompetensi anda {deskripsi_dirkom}.
        # Tugas anda adalah membuat judul materi yang sesuai dengan kompetensi yang akan diajarkan kepada peserta pelatihan. Judul materi wajib mencakup kompetensi sebagai berikut :
        # -----------------------------------------
        # {kompetensi}
        # -----------------------------------------
        # Tuliskan judul modul ajar pada format di bawah ini :
        # Judul :
        # """

        prompt = f"""system:
        Anda adalah ahli {dirkom} yang memiliki pengetahuan tentang {dirkom} dan memiliki kompetensi dalam {kompetensi}. 
        Anda diminta untuk membuat judul modul ajar yang komprehensif dan sesuai dengan kompetensi yang anda miliki.
        Anda tidak wajib memenuhi semua kompetensi yang ada, namun pastikan judul modul ajar yang anda buat sesuai dengan salah satu kompetensi yang anda miliki.
        user:Buatkanlah judul modul ajar yang sesuai dengan kompetensi yang saya miliki.
        Tuliskan judul modul ajar pada format di bawah ini :
        Judul : 
        assistant:"""
        return prompt

    def get_module_title(self, dirkom, kompetensi):
        title = self.request_generate(self.prompt_generate_title(dirkom, kompetensi))
        title = str(title)
        print(title)
        return title

    def prompt_generate_contentlist(self, title, dirkom, kompetensi):
        prompt = f"""system:
        Anda adalah ahli {dirkom} yang ahli dalam membuat modul ajar dan memiliki salah satu kompetensi dalam {kompetensi}. 
        Anda diminta untuk membuat daftar isi dari modul yang berjudul {title} berdasarkan dari referensi yang tersedia pada indeks yang tersimpan, BUAT JUDUL MAKSIMAL 8 KATA
        Daftar isi terdiri dari topik dan subtopik yang akan diajarkan. jumlah topik wajib MAXIMAL 3 topik dan julah subtopik wajib MINIMAL 2 subtopik per topik.
        user:Buatkanlah judul modul ajar yang sesuai dengan salah satu kompetensi yang anda miliki.
        Tuliskan daftar isi pada format di bawah ini :
        Judul : {title}
        topik 1 : 
            subtopik 1.1 :
            subtopik 1.2 :
            dan seterusnya
        topik 2 :
            subtopik 2.1 :
            subtopik 2.2 :
            dan seterusnya
        topik 3 :
            subtopik 3.1 :
            dan seterusnya
        
        assistant:
        """

        # prompt = f"""Anda adalah seorang ahli {dirkom} yang ahli dalam membuat modul ajar. kompetensi anda {deskripsi_dirkom}.
        # Tugas anda adalah membuat daftar isi dari modul yang berjudul {title} berdasarkan dari refrensi yang tersedia pada index yang disimpan.
        # Daftar isi terdiri dari topik dan subtopik yang akan diajarkan. jumlah topik wajib maksimal 3 topik dan jumlah subtopik wajib 2 subtopik per topik.
        # berikut adalah kompetensi yang wajib dicakup dalam modul ajar :
        # -----------------------------------------
        # {kompetensi}
        # -----------------------------------------
        # Tuliskan daftar isi pada format di bawah ini :
        # judul : {title}
        #     topik 1 :
        #         subtopik 1.1 :
        #         subtopik 1.2 :
        #         dan seterusnya
        #     topik 2 :
        #         subtopik 2.1 :
        #         subtopik 2.2 :
        #         dan seterusnya
        #     topik 3 :
        #         subtopik 3.1 :
        #         subtopik 3.2 :
        #         dan seterusnya
        # """
        return prompt

    def get_list_content(self, dirkom, kompetensi):
        title = self.get_module_title(dirkom, kompetensi)
        prompt = self.prompt_generate_contentlist(title, dirkom, kompetensi)
        list_content = self.request_generate(prompt)
        list_content = str(list_content)
        list_content = list_content.replace("\n", "").strip()
        print(list_content)
        topics = re.findall(
            r"Topik\s*(\d+)\s*:\s*(.+?)(?=topik|\Z|subtopik)",
            list_content,
            re.IGNORECASE,
        )
        subtopics = re.findall(
            r"Subtopik\s*(\d+\.\d+)\s*:\s*(.+?)(?=subtopik|\Z|topik)",
            list_content,
            re.IGNORECASE,
        )

        topics_dict = {topic[0]: {"topic": topic[1].strip()} for topic in topics}
        subtopics_dict = {
            subtopic[0]: {"subtopic": subtopic[1].strip()} for subtopic in subtopics
        }

        # Convert the dictionaries to JSON format
        result = {
            "judul": title.replace("Judul: ", ""),
            "Topic": topics_dict,
            "Subtopic": subtopics_dict,
        }
        return result

    def generate_content(self, dirkom, deskripsi_dirkom, kompetensi):
        content_list = self.get_list_content(dirkom, kompetensi)
        summary = ""
        module = {
            "judul": content_list["judul"],
            "topik": [],  # berisi dict subtopic (array)
            "references": [],
        }
        for topic in content_list["Topic"].items():
            module["topik"].append({"topic": topic[1], "subtopic": []})
            for subtopic in content_list["Subtopic"].items():
                if subtopic[0].startswith(topic[0]):
                    prompt = f"""system:
                    Anda adalah ahli {dirkom} yang ahli dalam membuat modul ajar dan memiliki kompetensi dalam {kompetensi}. 
                    Tugas anda adalah membuat konten untuk subtopik {subtopic[1]} dari topik {topic[1]}.
                    Pada perintah sebelumnya, anda telah membuat konten pada subtopik sebelumnya, berikut adalah rangkuman dari konten yang sebelumnya telah anda buat :
                    {summary}
                    Tambahkan pengetahuan kamu untuk membuat konten cocok untuk untuk metode ajar. Konten yang kamu buat harus terdiri dari 2 paragraf dan setiap paragraf terdiri dari 3 kalimat. 
                    Usahakan untuk tidak melebihi 70 kata. cantumkan referensi dari index yang disediakan dan gunakan in text reference dengan style APA.
                    Buatkan juga rangkuman singkat dari konten yang telah anda buat.
                    UNTUK IN TEXT REFERENCE JANGAN GUNAKAN NAMA FILE PDF YANG KAMU MASUKKAN KE TEKS, GUNAKAN TAHUN DAN AUTHORNYA SAJA. JIKA REFERENSINYA MELEBIHI 1 AUTHOR, GUNAKAN 'et.al'
                    Untuk summary pastikan kamu membuatnya maksimal 25 kata, pastikan untuk tidak melebihinya
                    user:Buatlah konten untuk subtopik {subtopic[1]} dari topik {topic[1]} sesuai dengan instruksi dari system.
                    Tuliskan konten menggunakan bahasa indonesia pada format di bawah ini :
                    konten:
                    paragraf 1:
                    paragraf 2:
                    rangkuman :
                    assistant:
                    """
                    # prompt = f"""Anda adalah seorang ahli {dirkom} yang ahli dalam membuat modul ajar. kompetensi anda {deskripsi_dirkom}.
                    # Tugas anda adalah membuat konten untuk subtopik {subtopic[1]} dari topik {topic[1]} dengan menggunakan referensi berikut : {references}.
                    # konten terdiri dari 2 paragraf dan setiap paragraf terdiri dari 3 kalimat. usahakan untuk tidak melebihi 50 kata. cantumkan referensi dari index yang disediakan dan gunakan in text reference.
                    # tulislah konten menggunakan bahasa indonesia pada format di bawah ini :
                    # konten :
                    # """
                    print(f"generate subtopic {subtopic[0]}")
                    content = self.request_generate(prompt)
                    result_tuple = list(content.metadata.items())[0]
                    result_dict = result_tuple[1]
                    publication_year = result_dict.get("publication_year")
                    author = result_dict.get("author")
                    file_name = result_dict.get("file_name")
                    references = f"{author}, {publication_year}. {file_name}"
                    content = str(content)
                    content = content.replace("\n", "").strip()
                    print(f"raw content: {content}")

                    # remove tulisan konten
                    content_text = re.sub(
                        r"\bkonten:\s*", "", content, flags=re.IGNORECASE
                    )

                    # remove tulisan summary
                    content_text = re.sub(
                        r"Rangkuman:.*|Rangkuman\s*:.",
                        "",
                        content,
                        flags=re.IGNORECASE | re.DOTALL,
                    )

                    # cari konten asli
                    content_text = re.findall(
                        r"(?:Paragraf\s*\d+\s*:?|Paragraf\d+:\s*|Paragraf\s*\d+\s*:\s*)(.*?)(?=\s*Paragraf\s*\d+\s*:?|Paragraf\d+:\s*|Paragraf\s*\d+\s*:\s*|$)",
                        content_text,
                        re.IGNORECASE,
                    )

                    summary = re.findall(
                        r"rangkuman\s*:\s*(.*?)(?=rangkuman|\Z)", content, re.IGNORECASE
                    )
                    # content_text = re.findall(r"Paragraf\s*\d\s*:\s*(.*?)\s*(?=Paragraf|\Z)", content_text, re.IGNORECASE)
                    print(topic[1])
                    print(subtopic[1])
                    print(content)
                    print(content_text)
                    print("SUMMARY: ", summary)
                    module["topik"][-1]["subtopic"].append(
                        {
                            "subtopic": subtopic[1],
                            "content": content_text,
                            "summary": summary,
                        }
                    )
                    if references not in module["references"]:
                        module["references"].append(references)
        print(module)

        import json

        with open("module.json", "w") as f:
            json.dump(module, f)

        return module

    def generate_summarization(self, topics):
        summaries = []

        for topic in topics:
            for subtopic in topic["subtopic"]:
                summaries.append(subtopic["summary"][0])
        prompt = f"""Kamu telah membuat content di perintah sebelumnya sebelumnya, tugas anda adalah generate ulang summary dari content yang anda buat dengan maksimal 15 kata setiap poin summary. atau maksimal total 120 kata jangan melebihinya 
        buat kedalam list poin poin penting dalam content
        berikut adalah content yang kamu generate sebelumnya:
        ---------------------------------------------------
        {summaries}
        ---------------------------------------------------
        Tuliskan summary pada format di bawah ini :
        - 

        - 

        - 

        - 

        - 

        - 

        - 
        --------------------------------------------------
        pastikan untuk menggunakan '-' pada setiap hasil summary, dan pastikan memberikan endline setelah menulis satu poin summaary
        """
        client = AzureOpenAI(
            api_key=os.getenv("AZURE_API_KEY"),
            api_version="2024-02-01",
            azure_endpoint=os.getenv("AZURE_API_BASE"),
        )
        print("Requesting completion...")
        deployment_name = "corpu-text-davinci-003"

        response = client.completions.create(
            model=deployment_name, prompt=prompt, temperature=0.4, max_tokens=2000
        )

        text = response.choices[0].text.replace(" .", ".").strip()
        return text

    """
    GENERATE SLIDE PRESENTATION
    """

    def preprocessing(self, text, max_length=100):
        # Replace special characters and split text into lines
        text = text.replace("\d", "").replace("\d. ", "").replace(". ", "\n").strip()
        text = re.sub(r"(^|\.\s)\d+\.\s", "\n", text)

        # Split text into lines with a maximum length of max_length
        words = text.split()
        lines = []
        current_line = ""
        for word in words:
            if len(current_line + word) <= max_length:
                current_line += word + " "
            else:
                lines.append(current_line.strip())
                current_line = word + " "
        if current_line:
            lines.append(current_line.strip())
        return "\n".join(lines)

    def get_date(self):
        bulan = {
            1: "Januari",
            2: "Februari",
            3: "Maret",
            4: "April",
            5: "Mei",
            6: "Juni",
            7: "Juli",
            8: "Agustus",
            9: "September",
            10: "Oktober",
            11: "November",
            12: "Desember",
        }

        now = datetime.now()
        date_display = "Last Update:" + now.strftime("%d {} %Y").format(
            bulan[now.month]
        )
        return date_display

    def add_picture_slide(self, slide, image_path, placeholder_idx):
        content_shape = slide.placeholders[placeholder_idx]
        content_shape.insert_picture(image_path)

    def add_slide_number(self, slide, prs):
        slide_number = prs.slides.index(slide) + 1

        x, y, width, height = Inches(12.61), Inches(6.5), Inches(0.6), Inches(0.6)
        text_box = slide.shapes.add_textbox(x, y, width, height)
        text_frame = text_box.text_frame

        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = str(slide_number)
        font = run.font
        font.name = "Arial"
        font.size = Pt(18)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)

    def add_title_slide(self, prs, title, image_path=None):
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_shape = title_slide.shapes.title
        subtitle_shape = title_slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = self.get_date()

        if image_path:
            self.add_picture_slide(title_slide, image_path, 2)

    def add_table_of_contents_slide(self, prs, data):
        # Pilih layout slide yang sesuai untuk daftar isi
        toc_slide_layout = prs.slide_layouts[1]
        toc_slide = prs.slides.add_slide(toc_slide_layout)

        title_placeholder = toc_slide.placeholders[0]
        title_placeholder.text = "Daftar Isi"

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        content_shape = toc_slide.shapes.add_table(
            rows=1, cols=2, left=left, top=top, width=width, height=height
        )
        table = content_shape.table

        table.cell(0, 0).text = "Topik"
        table.cell(0, 1).text = "Halaman"

        row_index = 1
        for topic in data["topik"]:
            for subtopic in topic["subtopic"]:
                table.cell(row_index, 0).text = subtopic["subtopic"]
                table.cell(row_index, 1).text = (
                    f'{prs.slides.index(subtopic["subtopic"]) + 1}'
                )
                row_index += 1

    def add_table_of_contents_slide(self, prs, data):
        toc_slide_layout = prs.slide_layouts[1]
        toc_slide = prs.slides.add_slide(toc_slide_layout)
        title_placeholder = toc_slide.placeholders[0]
        title_placeholder.text = "Daftar Isi"

        left, top, width, height = Inches(1), Inches(1.4), Inches(8), Inches(2.5)
        content_shape = toc_slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_shape.text_frame

        p = content_frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(14)

        for topic in data["topik"]:
            p = content_frame.add_paragraph()
            p.text = topic["topic"]["topic"] + "\n"
            p.level = 1
            # Ubah warna teks menjadi biru
            for run in p.runs:
                run.font.color.rgb = RGBColor(0, 0, 255)

            for subtopic in topic["subtopic"]:
                p = content_frame.add_paragraph()
                p.text = f'-    {subtopic["subtopic"]["subtopic"]}'
                p.level = 2
                p.space_after = Pt(7)

    def add_text_slide(
        self, slide, title, content, content_placeholder_index=1, line_spacing=1.5
    ):
        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title

        content_placeholder = slide.placeholders[content_placeholder_index]
        content_frame = content_placeholder.text_frame

        content = self.preprocessing(content)
        content_frame.text = content.replace("\n", ". ")

        # Ensure text is justified
        for paragraph in content_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.JUSTIFY
            paragraph.line_spacing = line_spacing

        content_frame.add_paragraph()

    def add_content_slide(self, prs, data, image_path=None):
        slide_layout = prs.slide_layouts[-1]

        # Create slide for the main content
        slide = prs.slides.add_slide(slide_layout)

        footer = slide.placeholders[1]

        if image_path:
            self.add_picture_slide(slide, image_path, 11)

        for topic in data["topik"]:
            topic_layout = prs.slide_layouts[-2]
            topic_slide = prs.slides.add_slide(topic_layout)
            self.add_text_slide(
                topic_slide, topic["topic"]["topic"].replace(": ", ""), ""
            )
            footer = topic_slide.placeholders[1]
            footer.text = f""
            if image_path:
                self.add_picture_slide(topic_slide, image_path, 2)

            for _, subtopic in enumerate(topic["subtopic"]):
                for i, content in enumerate(subtopic["content"]):
                    slide_content = prs.slides.add_slide(slide_layout)
                    self.add_text_slide(
                        slide_content,
                        subtopic["subtopic"]["subtopic"]
                        + f" ({i+1} / {len(subtopic['content'])})",
                        content,
                        content_placeholder_index=10,
                    )
                    footer = slide_content.placeholders[1]
                    footer.text = f""
                    if image_path:
                        self.add_picture_slide(slide_content, image_path, 11)

    def add_summary_slide(self, prs, summary):
        summary_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(summary_slide_layout)
        title_placeholder = summary_slide.placeholders[0]
        title_placeholder.text = "Ringkasan Konten"

        left, top, width, height = Inches(1), Inches(2.2), Inches(8), Inches(2.5)
        content_shape = summary_slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_shape.text_frame

        p = content_frame.add_paragraph()

        p = content_frame.add_paragraph()
        p.text = (
            summary.replace("Summary: ", "")
            .replace("Summary:", "")
            .replace("summary: ", "")
            .replace("summary:", "")
            .replace(". ", "\n")
            .replace(".", "\n")
        )
        p.space_after = Inches(0.2)
        p.level = 0

    def add_references_slide(self, prs, references):
        references_slide_layout = prs.slide_layouts[1]
        references_slide = prs.slides.add_slide(references_slide_layout)
        title_placeholder = references_slide.placeholders[0]
        title_placeholder.text = "Referensi"

        left, top, width, height = Inches(1), Inches(2.2), Inches(8), Inches(2.5)
        content_shape = references_slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_shape.text_frame

        # Tambahkan daftar referensi
        p = content_frame.add_paragraph()

        for reference in references:
            p = content_frame.add_paragraph()
            p.text = reference.replace("Referensi:", "").replace("\n", "")
            p.space_after = Inches(0.2)
            p.level = 0

    def generate_ppt(self, data, summary, output_path, image_path=None):
        prs = Presentation("template.pptx")
        # Add title slide
        self.add_title_slide(prs, data["judul"], image_path)

        self.add_table_of_contents_slide(prs, data)

        # Add content slides
        self.add_content_slide(prs, data, image_path)

        self.add_summary_slide(prs, summary)

        # Add references slide
        self.add_references_slide(prs, data["references"])

        del prs.slides._sldIdLst[0]
        del prs.slides._sldIdLst[0]
        del prs.slides._sldIdLst[2]

        # Add slide numbers
        for slide in prs.slides:
            self.add_slide_number(slide, prs)

        prs.save(output_path)
        return "Your ppt is ready"

    def clean(self, filename):
        cleaned_filename = re.sub(r'[<>:"/\\|?*\n]', "_", filename)
        return cleaned_filename

    def main(self, dirkom, deskripsi_dirkom, kompetensi):
        load_dotenv()
        data = self.generate_content(dirkom, deskripsi_dirkom, kompetensi)
        summary = self.generate_summarization(data["topik"])
        output_path = f"{dirkom}_l4_{data['judul']}.pptx"
        # output_path = "hasil" + dirkom + "_lx_" + data['judul'] + ".pptx"
        # output_path = "ppt_generated.pptx"
        # image_path = "business-plan_7788444.png"
        self.generate_ppt(data, summary, self.clean(output_path).replace(" ", "_"))
        print("success")


if __name__ == "__main__":
    start = time.time()

    fundamental = """
    Memiliki pengertian terhadap kompetensi terkait dan mampu mengingat pengetahuan faktual dan prosedural 
    Tidak memiliki atau memiliki sedikit pengalaman terhadap implementasi kompetensi
    Level ini biasanya diperlukan untuk posisi orang awam ke entry-level sehingga memerluka fundamental dalam kompetensi
    Modul ajar yang dibuat biasanya berisi pengetahuan dasar dan pengenalan terhadap kompetensi
    """

    developing = """
    Menunjukan pengetahuan dasar terkait kompetensi 
    Mulai menerapkan pengetahuan dalam pekerjaan namun membutuhkan bimbingan dalam implementasinya 
    Tahu bagaimana dan kapan harus mengeskalasi masalah
    Level ini biasanya diperlukan untuk posisi entry-level ke menengah sehingga memerlukan pengembangan dalam kompetensi
    Modul ajar yang dibuat biasanya berisi pengetahuan dasar dan bagaimana menerapkan kompetensinya
    """

    intermediate = """
    Menunjukan pengetahuan yang memadai dan mampu menjelaskan konsep-konsep kunci serta keterampilannya kepada orang lain 
    Mampu menerapkan konsep dan menggunakan pengetahuannya dalam aktivitas sehari-hari secara mandiri
    Mampu melakukan pemecahan masalah yang sederhana 
    Level ini biasanya diperlukan untuk posisi menengah atau bukan pemula ke hampir senior sehingga memerlukan pengembangan dalam kompetensi dan sudah fasih dalam hal dasar kompetensi
    Modul ajar yang dibuat sudah mengarah ke penerapan kompetensi dalam situasi nyata dan hampir mendekati level senior sehingga tidak diperlukan pembahasan yang terlalu dasar
    """

    advanced = """
    Menunjukan pengetahuan yang mendalam dan secara rutin mengaplikasikan konsep dan keterampilannya dalam situasi sulit
    Mampu membuat penilaian untuk menyelesaikan permasalahan rumit dan membuat keputusan penting terkait dengan impelementasi kompetensi 
    Mampu melatih orang lain dalam penerapan kompetensi
    Level ini biasanya diperlukan untuk posisi senior mampu mengajarkan kompetensi kepada orang lain dan mampu memberikan solusi terhadap permasalahan yang kompleks
    Modul ajar yang dibuat biasanya berisi penerapan kompetensi dalam situasi yang kompleks dan mendekati level master sehingga tidak diperlukan pembahasan yang terlalu dasar
    """

    mastery = """
    Diakui sebagai ahli yang memiliki pengetahuan, pemahaman dan penerapan kompetensi secara komperehensif 
    Menerapkan kompetensi dan mampu memecahkan permasalahan yang sangat kompleks di berbagi situasi 
    Mampu berinovasi, memformulasikan ide dan mengembangkan penerapan dan pedoman baru lintas organisasi
    Level ini biasanya diperlukan untuk posisi senior yang memiliki pengalaman yang sangat luas, dapat memahami kata teknis, mampu memberikan solusi terhadap permasalahan yang kompleks, dan menciptakan inovasi
    Modul ajar yang dibuat biasanya berisi penerapan kompetensi dalam situasi yang sangat kompleks dan dapat menciptakan inovasi sehingga tidak diperlukan pembahasan yang terlalu dasar
    """

    """
    boundary
    """

    dirkom = "Business Performance Management"

    deskripsi_dirkom = """
    
    """
    kompetensi_1 = (
        fundamental
        + """
    *Memiliki pemahaman dasar terkait manajemen performa bisnis termasuk pemahaman terkait KPI/OKR yang ditetapkan, prioritas bisnis saat ini, membaca dan memahami laporan keuangan dasar, serta proses evaluasi kinerja, perusahaan yang mencakup kinerja portfolio, operasional, serta finansial perusahaan dan/atau anak perusahaan
    *Memiliki pengetahuan terkait metode analisis pengukurang performa perusahaan, meliputi namun tidak terbatas pada balanced sVcard, strategic planning and management system, performance management cycle, serta metrik penilaian yang digunakan dalam menilai performa perusahaan dan/atau anak perusahaan (Seperti key performance indicators, leading indicators, dan lagging indicators) 
    """
    )
    kompetensi_2 = (
        developing
        + """
    *Memiliki pemahaman terkait laporan keuangan lanjutan seperti laporan pajak, neraca, arus kas, dan sebagainya
    *Mampu mengidentifikasi metodologi yang tepat untuk digunakan dalam proses penilaian performa perusahaan dan/atau anak perusahaan 
    *Mampu melakukan pengumpulan data terkait performa perusahaan, yaitu data internal (seperti data keuangan, operasonal, dan non-finansial) dan data eksternal (seperti indikator ekonomi, data industri telekomunikasi dan digital, dan respon pada media sosial)  untuk melakukan identifikasi terhadap gap kondisi aktual dan ideal perusahaan dan/atau anak perusahaan
    """
    )
    kompetensi_3 = (
        intermediate
        + """
    *Mampu melakukan penilaian terhadap performa menggunakan metode penilaian performa yang telah ditetapkan 
    *Mampu melakukan analisis terkait gap terkait performa aktual dan ideal perusahaan untuk mengidentifikasi faktor penyebab dari  tidak tercapainya indikator performa perusahaan dan/atau anak perusahaan dalam rangka perancangan laporan performa perusahaan secara menyeluruh
    """
    )
    kompetensi_4 = (
        advanced
        + """
    *Mampu melakukan evaluasi terhadap performa perusahaan dan/atau anak perusahaan berdasarkan hasil analisis dan pengolahan data perofrma dan mengembangkan strategi perbaikan pada area yang membutuhkan perbaikan 
    *Mampu melakukan evaluasi terkait KPI/OKR yang berlaku untuk melakukan optimalisasi batas ketercapaian performa dan memastikan KPI/OKR mewakilkan setiap proses bisnis perusahaan
    *Mampu mengevaluasi ketercapaian performa berdasarkan target yang telah ditetapkan dan memberikan strategi perbaikan, dan merancang strategi optimalisasi performa bisnis yang mencakup kinerja portfolio, operasional, dan finansial 
    *Mampu melakukan evaluasi terhadap performa bisnis termasuk kinerja portfolio, operasional, dan finansial untuk mengidentifikasi kekurangan dan area perbaikan dan memberika rekomendasi strategi optimalisasi atau pengembangan 

    """
    )
    kompetensi_5 = (
        mastery
        + """
    *Mampu memberikan rekomendasi strategis untuk proses peningkatan performa perusahaan secara holistik  dengan mempertimbangkan perkembangan tren pasar, kebijakan publik, serta perubahan pada persaingan industri untuk meningkatkan efektivitas dan efisiensi
    *Mampu memberikan solusi yang komprehensif terkait kendala pada program penjagaan atau peningkatan performa maupun langkah antisipatif terkait turunnya performa pekerja melalui inovasi pada pendekatan dan metode analisis performa pekerja 
    *Mampu memberikan rekomendasi atau saran terhadap strategi dan implementasi manajemen performa bisnis termasuk prosedur penilaian performa yang ditetapkan, program pengembangan yang dilakukan, serta perancangan KPI/OKR atau metrik pengukuran performa yang digunakan dalam perusahaan dan/atau anak perusahaan
    *Mampu melakukan analisis dampak jangka panjang yang diakibatkan jika tidak dilakukannya perbaikan pada performa proses bisnis yang tidak memenuhi KPI/OKR
    *Mampu melakukan analisis dan proyeksi terhadap perkembangan teknologi yang relevan dengan pemukuran performa untuk melakukan inovasi pada proses pengukura hingga pelaporan performa perusahaan melalui implementasi teknologi digital dalam rangka meningkatkan efektivitas dan efisiensi
    """
    )
    generator = ContentGenerator()
    generator.main(dirkom, deskripsi_dirkom, kompetensi_1)
    # generator.main(dirkom, deskripsi_dirkom, kompetensi_2)
    # generator.main(dirkom, deskripsi_dirkom, kompetensi_3)
    # generator.main(dirkom, deskripsi_dirkom, kompetensi_4)
    # generator.main(dirkom, deskripsi_dirkom, kompetensi_5)
    end = time.time()
    elapsed = end - start
    print(f"time elapsed: {elapsed}")
